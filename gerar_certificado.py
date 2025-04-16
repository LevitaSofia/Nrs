import sqlite3
import pandas as pd
from pptx import Presentation
import os
from datetime import datetime
import shutil

def criar_modelos_nrs():
    """Cria modelos para todas as NRs baseado no modelo NR06"""
    modelo_base = "Certificado - NR06.pptx"
    if not os.path.exists(modelo_base):
        print(f"Erro: Modelo base {modelo_base} não encontrado!")
        return False
    
    # Criar diretório para modelos se não existir
    if not os.path.exists("modelos"):
        os.makedirs("modelos")
    
    # Informações de cada NR
    nrs_info = {
        5: {
            "titulo": "NR 05",
            "nome_curso": "Curso de CIPA - Comissão Interna de Prevenção de Acidentes",
            "carga_horaria": "20"
        },
        6: {
            "titulo": "NR 06",
            "nome_curso": "Curso sobre uso e guarda de EPI",
            "carga_horaria": "04"
        },
        10: {
            "titulo": "NR 10",
            "nome_curso": "Curso de Segurança em Instalações e Serviços com Eletricidade",
            "carga_horaria": "40"
        },
        11: {
            "titulo": "NR 11",
            "nome_curso": "Curso de Operação de Equipamentos de Transporte de Cargas",
            "carga_horaria": "16"
        },
        12: {
            "titulo": "NR 12",
            "nome_curso": "Curso de Segurança no Trabalho em Máquinas e Equipamentos",
            "carga_horaria": "08"
        },
        17: {
            "titulo": "NR 17",
            "nome_curso": "Curso de Ergonomia",
            "carga_horaria": "08"
        },
        20: {
            "titulo": "NR 20",
            "nome_curso": "Curso de Segurança com Inflamáveis e Combustíveis",
            "carga_horaria": "16"
        },
        33: {
            "titulo": "NR 33",
            "nome_curso": "Curso de Segurança e Saúde nos Trabalhos em Espaços Confinados",
            "carga_horaria": "16"
        },
        35: {
            "titulo": "NR 35",
            "nome_curso": "Curso de Trabalho em Altura",
            "carga_horaria": "08"
        }
    }
    
    # Criar modelo para cada NR
    for nr, info in nrs_info.items():
        novo_modelo = f"modelos/Certificado - NR{nr:02d}.pptx"
        # Copiar o modelo base
        shutil.copy2(modelo_base, novo_modelo)
        
        # Abrir o novo modelo para edição
        prs = Presentation(novo_modelo)
        
        # Para cada slide no modelo
        for slide in prs.slides:
            # Para cada shape no slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Substituir textos mantendo a formatação
                    texto = shape.text
                    texto = texto.replace("NR 06", info["titulo"])
                    texto = texto.replace("Curso sobre uso e guarda de EPI", info["nome_curso"])
                    texto = texto.replace("04 horas", f"{info['carga_horaria']} horas")
                    shape.text = texto
        
        # Salvar o modelo
        prs.save(novo_modelo)
    
    return True

def listar_colaboradores():
    conn = sqlite3.connect("colaboradores.db")
    df = pd.read_sql("SELECT rowid, Nome, CPF FROM colaboradores", conn)
    conn.close()
    return df

def selecionar_colaborador():
    df = listar_colaboradores()
    print("\nLista de Colaboradores:")
    for idx, row in df.iterrows():
        print(f"{row['rowid']}. {row['Nome']} - CPF: {row['CPF']}")
    
    while True:
        try:
            escolha = int(input("\nDigite o número do colaborador: "))
            if escolha in df['rowid'].values:
                return df[df['rowid'] == escolha].iloc[0]
            else:
                print("Número inválido. Tente novamente.")
        except ValueError:
            print("Por favor, digite um número válido.")

def selecionar_nrs():
    nrs_info = {
        1: {
            "numero": "06",
            "titulo": "NR 06",
            "nome_curso": "Curso sobre uso e guarda de EPI",
            "carga_horaria": "04",
            "modelo": "Certificado - NR06.pptx"
        },
        2: {
            "numero": "12",
            "titulo": "NR 12",
            "nome_curso": "Curso de Segurança no Trabalho em Máquinas e Equipamentos",
            "carga_horaria": "08",
            "modelo": "Certificado - NR12.pptx"
        },
        3: {
            "numero": "18",
            "titulo": "NR 18",
            "nome_curso": "Curso de Segurança na Construção Civil",
            "carga_horaria": "16",
            "modelo": "Certificado - NR18.pptx"
        },
        4: {
            "numero": "35",
            "titulo": "NR 35",
            "nome_curso": "Curso de Trabalho em Altura",
            "carga_horaria": "08",
            "modelo": "Certificado - NR35.pptx"
        }
    }
    
    print("\nTipos de NRS disponíveis:")
    for key, info in nrs_info.items():
        print(f"{key}. {info['titulo']} - {info['nome_curso']}")
    
    while True:
        try:
            escolha = int(input("\nDigite o número da NR: "))
            if escolha in nrs_info:
                return nrs_info[escolha]
            else:
                print("Número inválido. Tente novamente.")
        except ValueError:
            print("Por favor, digite um número válido.")

def substituir_texto_shape(shape, texto_antigo, texto_novo):
    """Substitui texto mantendo a formatação básica"""
    if not hasattr(shape, "text"):
        return
        
    if texto_antigo in shape.text:
        # Guardar formatação dos parágrafos
        paragraph_formats = []
        run_formats = []
        
        for paragraph in shape.text_frame.paragraphs:
            # Guardar formatação do parágrafo
            p_format = {
                'alignment': paragraph.alignment,
                'level': paragraph.level
            }
            paragraph_formats.append(p_format)
            
            # Guardar formatação dos runs
            runs_format = []
            for run in paragraph.runs:
                r_format = {
                    'font_name': run.font.name,
                    'size': run.font.size,
                    'bold': run.font.bold,
                    'italic': run.font.italic
                }
                runs_format.append(r_format)
            run_formats.append(runs_format)
        
        # Fazer a substituição
        shape.text = shape.text.replace(texto_antigo, texto_novo)
        
        # Restaurar formatação
        for i, paragraph in enumerate(shape.text_frame.paragraphs):
            if i < len(paragraph_formats):
                p_format = paragraph_formats[i]
                paragraph.alignment = p_format['alignment']
                paragraph.level = p_format['level']
                
                # Restaurar formatação dos runs
                if i < len(run_formats):
                    for j, run in enumerate(paragraph.runs):
                        if j < len(run_formats[i]):
                            r_format = run_formats[i][j]
                            if r_format['font_name']:
                                run.font.name = r_format['font_name']
                            if r_format['size']:
                                run.font.size = r_format['size']
                            run.font.bold = r_format['bold']
                            run.font.italic = r_format['italic']

def gerar_certificado(colaborador, nr_info):
    # Verificar se o modelo existe
    modelo_arquivo = nr_info['modelo']
    if not os.path.exists(modelo_arquivo):
        print(f"Erro: Modelo {modelo_arquivo} não encontrado!")
        return None
    
    # Criar diretório para certificados se não existir
    if not os.path.exists("certificados"):
        os.makedirs("certificados")
    
    # Nome do novo certificado
    nome_arquivo = f"certificados/Certificado_{colaborador['Nome'].replace(' ', '_')}_NR{nr_info['numero']}.pptx"
    
    # Copiar o modelo
    shutil.copy2(modelo_arquivo, nome_arquivo)
    
    # Abrir o novo certificado
    prs = Presentation(nome_arquivo)
    
    # Formatar a data de ministração
    data_ministerio = datetime.strptime(nr_info['data_ministerio'], '%Y-%m-%d').strftime('%d/%m/%Y')
    
    # Textos a serem substituídos
    substituicoes = [
        ("{{NOME}}", colaborador['Nome']),
        ("{{CPF}}", colaborador['CPF']),
        ("{{DATA}}", data_ministerio)  # Agora usa {{DATA}} em vez da data fixa
    ]
    
    # Para cada slide
    for slide in prs.slides:
        # Para cada shape no slide
        for shape in slide.shapes:
            # Para cada substituição necessária
            for texto_antigo, texto_novo in substituicoes:
                substituir_texto_shape(shape, texto_antigo, texto_novo)
    
    # Salvar o certificado
    prs.save(nome_arquivo)
    print(f"\nCertificado gerado com sucesso: {nome_arquivo}")
    
    return nome_arquivo

def main():
    print("Sistema de Geração de Certificados NRS")
    print("=====================================")
    
    colaborador = selecionar_colaborador()
    nr_info = selecionar_nrs()
    gerar_certificado(colaborador, nr_info)

if __name__ == "__main__":
    main() 