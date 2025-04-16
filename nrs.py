import pandas as pd
from pptx import Presentation
import os

# Arquivos
arquivo_excel = "lista_colaboradores.xlsx"
modelo_ppt = "nr06.pptx"
saida_dir = "certificados"

# Garantir a pasta de saída
os.makedirs(saida_dir, exist_ok=True)

# Carrega os dados da planilha
df = pd.read_excel(arquivo_excel)

# Função para substituir campos no slide
def substituir_campos(slide, dados):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for par in shape.text_frame.paragraphs:
                for run in par.runs:
                    for chave, valor in dados.items():
                        if chave in run.text:
                            run.text = run.text.replace(chave, str(valor))

# Para cada pessoa na planilha, gerar um slide
for _, row in df.iterrows():
    prs = Presentation(modelo_ppt)
    slide = prs.slides[0]
    
    dados = {
        "{{NOME}}": row["Nome"],
        "{{CPF}}": row["CPF"]
    }

    substituir_campos(slide, dados)

    nome_arquivo = f"{saida_dir}/Certificado_{row['Nome'].replace(' ', '_')}.pptx"
    prs.save(nome_arquivo)



print("Todos os certificados foram gerados!")
